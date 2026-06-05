import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def html_to_pptx(html_path, pptx_path):
    with open(html_path, 'r', encoding='utf-8') as f:
        html = f.read()

    prs = Presentation()
    
    # Extract slides
    slides_raw = re.findall(r'<section class="slide.*?".*?>(.*?)</section>', html, re.DOTALL)
    
    for slide_content in slides_raw:
        # Create a blank slide
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]
        
        # Extract title (h1 or h2)
        title_match = re.search(r'<h[12].*?>(.*?)</h[12]>', slide_content, re.DOTALL)
        if title_match:
            title_text = re.sub(r'<.*?>', '', title_match.group(1)).strip()
            title_placeholder.text = title_text
            
        # Extract body text (p and benefit-card content)
        body_text = []
        
        # Find paragraphs
        paragraphs = re.findall(r'<p.*?>(.*?)</p>', slide_content, re.DOTALL)
        for p in paragraphs:
            text = re.sub(r'<.*?>', '', p).strip()
            if text and not text.startswith('AI Automation'): # Skip footer
                body_text.append(text)
                
        # Find benefit cards or nodes
        nodes = re.findall(r'<div class="(?:node-title|benefit-card).*?".*?>(.*?)</div>', slide_content, re.DOTALL)
        for node in nodes:
            text = re.sub(r'<.*?>', '', node).strip()
            if text:
                body_text.append(text)
                
        # Fill body placeholder
        if body_text:
            tf = body_placeholder.text_frame
            tf.text = body_text[0]
            for i in range(1, len(body_text)):
                p = tf.add_paragraph()
                p.text = body_text[i]
                p.level = 0

    prs.save(pptx_path)
    print(f"Successfully saved to {pptx_path}")

if __name__ == "__main__":
    html_to_pptx('n8n-jupviec-automation.html', 'n8n-jupviec-automation.pptx')
